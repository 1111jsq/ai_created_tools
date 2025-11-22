from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from .models import PresentationRequest, ChartType, ThemeStyle

class PPTRenderer:
    def __init__(self, style_name: str = "tech"):
        self.theme = self._get_theme(style_name)

    def _get_theme(self, name: str) -> ThemeStyle:
        themes = {
            "tech": ThemeStyle(
                name="tech",
                bg_color=RGBColor(30, 30, 40), # Dark Blue-Grey
                text_color=RGBColor(240, 240, 240), # Off-white
                card_color=RGBColor(50, 50, 60), # Lighter dark
                accent_colors=[
                    RGBColor(0, 255, 255),    # Cyan
                    RGBColor(255, 0, 255),    # Magenta
                    RGBColor(255, 215, 0),    # Gold
                    RGBColor(50, 205, 50),    # Lime Green
                    RGBColor(255, 69, 0),     # Red Orange
                ]
            ),
            "light": ThemeStyle(
                name="light",
                bg_color=RGBColor(255, 255, 255), # White
                text_color=RGBColor(40, 40, 40),  # Dark Grey
                card_color=RGBColor(245, 245, 250), # Very light grey
                accent_colors=[
                    RGBColor(52, 152, 219),   # Blue
                    RGBColor(46, 204, 113),   # Green
                    RGBColor(155, 89, 182),   # Purple
                    RGBColor(231, 76, 60),    # Red
                    RGBColor(241, 196, 15),   # Yellow
                ],
                title_font="Arial",
                body_font="Calibri"
            ),
            "retro": ThemeStyle(
                name="retro",
                bg_color=RGBColor(253, 246, 227), # Solarized Base3
                text_color=RGBColor(88, 110, 117), # Solarized Base01
                card_color=RGBColor(238, 232, 213), # Solarized Base2
                accent_colors=[
                    RGBColor(211, 54, 130),   # Magenta
                    RGBColor(108, 113, 196),  # Violet
                    RGBColor(38, 139, 210),   # Blue
                    RGBColor(42, 161, 152),   # Cyan
                    RGBColor(133, 153, 0),    # Green
                ],
                title_font="Courier New",
                body_font="Courier New"
            )
        }
        return themes.get(name, themes["tech"])

    def create_presentation(self, request: PresentationRequest, output_file: str = "output.pptx"):
        prs = Presentation()
        
        # Configure Slide Dimensions (Widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Use a blank layout
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. Apply Background
        self._apply_background(slide)

        # 2. Add Decoration (Header Bar or shapes)
        self._add_decorations(slide)

        # 3. Add Title (Poster Style: Big, Bold)
        self._add_title(slide, request.slide_title)
        
        # 4. Add Chart in a "Card" container
        self._add_chart(slide, request)

        # 5. Add Summary in a "Card" container
        self._add_summary(slide, request.summary)

        prs.save(output_file)
        print(f"Presentation saved to {output_file} using theme '{self.theme.name}'")

    def _apply_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme.bg_color

    def _add_decorations(self, slide):
        # Add a subtle accent line at the top
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(13.333), Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme.accent_colors[0]
        shape.line.fill.background() # No border

    def _add_title(self, slide, text):
        # Left aligned, large modern font
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12)
        height = Inches(1.0)
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.theme.title_font
        p.font.size = Pt(44)
        # Use 3rd accent color for title if available, else 1st
        title_color_idx = 2 if len(self.theme.accent_colors) > 2 else 0
        p.font.color.rgb = self.theme.accent_colors[title_color_idx]

    def _add_chart(self, slide, request):
        # Layout metrics
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(8.5)
        height = Inches(5.0)

        # 1. Draw Card Background for Chart
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.theme.card_color
        # Conditional border
        if self.theme.name == "tech":
            bg_shape.line.color.rgb = RGBColor(80, 80, 90)
        else:
            bg_shape.line.fill.background()
        
        # 2. Prepare Data
        chart_data = CategoryChartData()
        chart_data.categories = request.data.categories
        for series in request.data.series:
            chart_data.add_series(series.name, series.values)

        # 3. Determine Type
        ppt_chart_type = self._get_ppt_chart_type(request.chart_type)

        # 4. Add Chart (Inset slightly from card)
        chart = slide.shapes.add_chart(
            ppt_chart_type, 
            left + Inches(0.2), top + Inches(0.2), 
            width - Inches(0.4), height - Inches(0.4), 
            chart_data
        ).chart

        # 5. Style Chart
        self._style_chart_fonts(chart)
        self._style_chart_colors(chart, request.chart_type)

    def _add_summary(self, slide, summary_text):
        # Right side column for insights
        left = Inches(9.3)
        top = Inches(1.8)
        width = Inches(3.5)
        height = Inches(5.0)

        # Card Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.theme.card_color
        bg_shape.line.fill.background()

        # "Key Insights" Header
        header_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
        hp = header_box.text_frame.paragraphs[0]
        hp.text = "KEY INSIGHTS"
        hp.font.name = self.theme.body_font
        hp.font.bold = True
        hp.font.size = Pt(14)
        hp.font.color.rgb = self.theme.accent_colors[1] # 2nd accent

        # Body Text
        body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(1.0))
        tf = body_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = summary_text
        p.font.name = self.theme.body_font
        p.font.size = Pt(16)
        p.font.color.rgb = self.theme.text_color

    def _get_ppt_chart_type(self, c_type):
        if c_type == ChartType.BAR: return XL_CHART_TYPE.COLUMN_CLUSTERED
        if c_type == ChartType.LINE: return XL_CHART_TYPE.LINE
        if c_type == ChartType.PIE: return XL_CHART_TYPE.DOUGHNUT
        if c_type == ChartType.SCATTER: return XL_CHART_TYPE.XY_SCATTER
        return XL_CHART_TYPE.COLUMN_CLUSTERED

    def _style_chart_fonts(self, chart):
        # Set font sizes and colors for axes and legends
        try:
            chart.font.name = self.theme.body_font
            chart.font.size = Pt(12)
            chart.font.color.rgb = self.theme.text_color
        except:
            pass
        
        # Legend
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = self.theme.text_color
            chart.legend.font.size = Pt(12)

        # Axis Labels
        try:
            chart.category_axis.tick_labels.font.color.rgb = self.theme.text_color
            chart.category_axis.tick_labels.font.size = Pt(11)
            chart.value_axis.tick_labels.font.color.rgb = self.theme.text_color
            chart.value_axis.tick_labels.font.size = Pt(11)
            
            # Gridlines color based on theme
            grid_color = RGBColor(70, 70, 80) if self.theme.name == "tech" else RGBColor(200, 200, 200)
            chart.value_axis.major_gridlines.format.line.color.rgb = grid_color
        except:
            pass

    def _style_chart_colors(self, chart, c_type):
        try:
            colors = self.theme.accent_colors
            for i, series in enumerate(chart.series):
                color = colors[i % len(colors)]
                
                # For line charts, color the line
                if c_type == ChartType.LINE:
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(3)
                # For bar/column, fill the shape
                elif c_type == ChartType.BAR:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                # For Pie/Doughnut, we need to color data points individually
                elif c_type == ChartType.PIE:
                    for j, point in enumerate(series.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = colors[j % len(colors)]
                elif c_type == ChartType.SCATTER:
                    series.marker.format.fill.solid()
                    series.marker.format.fill.fore_color.rgb = color
        except Exception as e:
            print(f"Warning: Could not style chart fully: {e}")
